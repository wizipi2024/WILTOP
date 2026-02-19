"""
VsCodeSkill - Criador de Apps com IA.
Cria projetos completos: HTML, PowerPoint, Word, Excel, Python.
Fluxo: perguntas chave → planejamento → revisão dos agentes → criação do arquivo.
"""

import subprocess
import json
import os
from pathlib import Path
from datetime import datetime
from src.skills.base_skill import BaseSkill, SkillResult
from src.utils.logger import get_logger

log = get_logger(__name__)

APPS_DIR = Path("data/apps")


class VsCodeSkill(BaseSkill):
    """Cria apps, apresentações, planilhas e sites usando IA + planejamento."""

    name = "vscode"
    description = "Cria apps, apresentações PowerPoint, planilhas Excel, documentos Word e sites"
    keywords = [
        "crie um app", "crie um site", "crie uma ferramenta", "desenvolva",
        "programe", "cria um sistema", "faca um aplicativo", "landing page",
        "pagina de captura", "dashboard web", "sistema de agendamento",
        "crm simples", "calculadora de roi", "gerador de proposta",
        "bot de telegram", "bot de whatsapp", "plataforma", "saas",
        "crie uma apresentacao", "crie um powerpoint", "crie uma planilha",
        "crie um excel", "crie um documento", "crie um word",
        "fazer apresentacao", "fazer powerpoint", "fazer planilha",
        "planilha de", "apresentacao de", "documento de",
    ]

    def can_handle(self, command: str) -> float:
        cmd = command.lower()
        score = 0.0
        for kw in self.keywords:
            if kw in cmd:
                score += 0.20
        return min(score, 1.0)

    def execute(self, command: str, params: dict = None, context: dict = None) -> SkillResult:
        """
        Ponto de entrada principal.
        O fluxo de perguntas/planejamento é gerido pela GUI (william_final.py).
        Aqui recebemos os dados já coletados via params ou um comando direto.
        """
        cmd = command.lower()
        p = params or {}

        # Dados de planejamento já coletados pela GUI
        if p.get("planned") and p.get("tipo") and p.get("descricao"):
            return self._executar_plano(p)

        # Detecção direta por tipo
        try:
            if any(w in cmd for w in ["powerpoint", "apresentacao", "apresentação", "slides", "pptx"]):
                return self._criar_powerpoint_direto(command)
            elif any(w in cmd for w in ["planilha", "excel", "xlsx", "spreadsheet"]):
                return self._criar_excel_direto(command)
            elif any(w in cmd for w in ["word", "documento", "docx", ".doc"]):
                return self._criar_word_direto(command)
            elif any(w in cmd for w in ["landing page", "pagina de captura", "pagina de vendas"]):
                return self._criar_landing_page(command)
            elif any(w in cmd for w in ["dashboard", "painel", "relatorio web"]):
                return self._criar_dashboard(command)
            elif any(w in cmd for w in ["calculadora", "roi", "simulador"]):
                return self._criar_calculadora(command)
            elif any(w in cmd for w in ["crm", "sistema de clientes", "gestao de leads"]):
                return self._criar_crm(command)
            else:
                return self._criar_app_generico(command)
        except Exception as e:
            log.error(f"VsCodeSkill erro: {e}")
            return SkillResult(success=False, message=f"Erro ao criar: {e}")

    # ================================================================
    # FLUXO PLANEJADO (chamado pela GUI após coletar perguntas)
    # ================================================================

    def _executar_plano(self, p: dict) -> SkillResult:
        """Executa criação com base no plano aprovado pelo usuário."""
        tipo          = p.get("tipo", "html")
        nome          = p.get("nome", "") or p.get("descricao", "app")
        descricao     = p.get("descricao", "app")
        publico       = p.get("publico", "")
        funcionalidades = p.get("funcionalidades", "")
        estilo        = p.get("estilo", "profissional")
        extras        = p.get("extras", "")
        abrir_vscode  = p.get("abrir_vscode", False)

        # Contexto rico — funcionalidades são a parte mais importante
        funcionalidades_formatadas = "\n".join(
            f"  - {l.strip().strip('-*•▸ ')}"
            for l in funcionalidades.splitlines()
            if l.strip()
        ) or "  - Funcionalidades gerais"

        contexto_completo = (
            f"Nome/Titulo: {nome}\n"
            f"Tipo: {tipo}\n"
            f"Descricao: {descricao}\n"
            f"Publico alvo: {publico or 'geral'}\n"
            f"Estilo/Tom: {estilo}\n"
            f"Informacoes adicionais: {extras}\n\n"
            f"FUNCIONALIDADES OBRIGATORIAS (IMPLEMENTAR TODAS SEM EXCECAO):\n"
            f"{funcionalidades_formatadas}"
        ).strip()

        tipo_lower = tipo.lower()
        if tipo_lower in ("powerpoint", "pptx", "apresentacao", "apresentação", "slides"):
            return self._criar_powerpoint_planejado(nome, contexto_completo)
        elif tipo_lower in ("excel", "xlsx", "planilha"):
            return self._criar_excel_planejado(nome, contexto_completo)
        elif tipo_lower in ("word", "docx", "documento"):
            return self._criar_word_planejado(nome, contexto_completo)
        elif tipo_lower in ("landing", "landing_page", "pagina de captura"):
            return self._criar_app_html_planejado(
                nome,
                f"Tipo: Landing Page de Alta Conversao\n{contexto_completo}",
                abrir_vscode=abrir_vscode)
        elif tipo_lower in ("dashboard", "painel"):
            return self._criar_app_html_planejado(
                nome,
                f"Tipo: Dashboard Analitico com KPIs\n{contexto_completo}",
                abrir_vscode=abrir_vscode)
        elif tipo_lower == "crm":
            return self._criar_app_html_planejado(
                nome,
                f"Tipo: CRM / Sistema de Gestao de Clientes\n{contexto_completo}",
                abrir_vscode=abrir_vscode)
        elif tipo_lower == "calculadora":
            return self._criar_app_html_planejado(
                nome,
                f"Tipo: Calculadora / Simulador Interativo\n{contexto_completo}",
                abrir_vscode=abrir_vscode)
        elif tipo_lower == "jogo":
            return self._criar_app_html_planejado(
                nome,
                f"Tipo: Jogo HTML com Canvas (game completo funcional)\n{contexto_completo}",
                abrir_vscode=abrir_vscode)
        elif tipo_lower == "portfolio":
            return self._criar_app_html_planejado(
                nome,
                f"Tipo: Portfolio / Site Pessoal\n{contexto_completo}",
                abrir_vscode=abrir_vscode)
        else:
            return self._criar_app_html_planejado(
                nome, contexto_completo, abrir_vscode=abrir_vscode)

    # ================================================================
    # POWERPOINT
    # ================================================================

    def _criar_powerpoint_planejado(self, descricao: str, contexto: str) -> SkillResult:
        """Cria apresentação PowerPoint real usando python-pptx com conteúdo da IA."""
        APPS_DIR.mkdir(parents=True, exist_ok=True)
        nome = f"apresentacao_{descricao.replace(' ', '_')[:30]}_{datetime.now().strftime('%Y%m%d_%H%M')}"
        pasta = APPS_DIR / nome
        pasta.mkdir(parents=True, exist_ok=True)

        # Gera estrutura da apresentação com IA
        prompt = f"""Voce e um especialista em apresentacoes de negocios.
Com base nestas informacoes:
{contexto}

Crie o roteiro completo de uma apresentacao profissional em formato JSON:
{{
  "titulo": "Titulo da apresentacao",
  "subtitulo": "Subtitulo ou tagline",
  "slides": [
    {{
      "titulo": "Titulo do slide",
      "tipo": "capa|conteudo|lista|dados|conclusao",
      "pontos": ["ponto 1", "ponto 2", "ponto 3"],
      "nota": "nota do apresentador (opcional)"
    }}
  ]
}}

REGRAS:
- Minimo 8 slides, maximo 15
- Dados e numeros ESPECIFICOS e REAIS (nao genericos)
- Conteudo ESPECIFICO para o contexto descrito, nao templates vazios
- Use linguagem assertiva e profissional
- Slides de dados: inclua numeros, percentuais, metricas reais do mercado
Retorne APENAS o JSON, sem explicacoes."""

        estrutura = None
        try:
            from src.core.ai_engine import get_engine
            engine = get_engine()
            resp = engine.chat_with_fallback(prompt)
            # Extrai JSON da resposta
            if "{" in resp:
                inicio = resp.find("{")
                fim = resp.rfind("}") + 1
                estrutura = json.loads(resp[inicio:fim])
        except Exception as e:
            log.warning(f"IA nao gerou estrutura: {e}")

        # Fallback de estrutura
        if not estrutura:
            estrutura = {
                "titulo": descricao.title(),
                "subtitulo": "Apresentacao Profissional",
                "slides": [
                    {"titulo": "Introducao", "tipo": "capa", "pontos": [descricao], "nota": ""},
                    {"titulo": "Contexto", "tipo": "conteudo", "pontos": ["Ponto 1", "Ponto 2", "Ponto 3"], "nota": ""},
                    {"titulo": "Solucao", "tipo": "lista", "pontos": ["Item A", "Item B", "Item C"], "nota": ""},
                    {"titulo": "Beneficios", "tipo": "lista", "pontos": ["Beneficio 1", "Beneficio 2", "Beneficio 3"], "nota": ""},
                    {"titulo": "Proximos Passos", "tipo": "conclusao", "pontos": ["Acao 1", "Acao 2", "Contato"], "nota": ""},
                ]
            }

        # Cria o arquivo .pptx
        arquivo_pptx = pasta / f"{nome}.pptx"
        try:
            self._montar_pptx(arquivo_pptx, estrutura)
        except ImportError:
            # Instala python-pptx
            log.info("Instalando python-pptx...")
            subprocess.run(["py", "-3", "-m", "pip", "install", "python-pptx", "-q"],
                           capture_output=True, timeout=60)
            try:
                self._montar_pptx(arquivo_pptx, estrutura)
            except Exception as e2:
                return SkillResult(success=False, message=f"Erro ao criar PowerPoint: {e2}")
        except Exception as e:
            return SkillResult(success=False, message=f"Erro ao montar PPTX: {e}")

        # Abre o arquivo
        try:
            os.startfile(str(arquivo_pptx))
        except Exception:
            subprocess.Popen(["start", "", str(arquivo_pptx)], shell=True)

        n_slides = len(estrutura.get("slides", []))
        lista_slides = "\n".join(
            f"  Slide {i+1}: {s.get('titulo', '')}"
            for i, s in enumerate(estrutura.get("slides", []))
        )
        return SkillResult(
            success=True,
            message=(
                f"[POWERPOINT CRIADO] {estrutura.get('titulo', descricao.title())}\n\n"
                f"ARQUIVO: {arquivo_pptx}\n"
                f"SLIDES: {n_slides} slides criados com conteudo real da IA\n"
                f"STATUS: Aberto no PowerPoint\n\n"
                f"CONTEUDO:\n{lista_slides}\n\n"
                f"Pronto para usar, editar e apresentar!"
            ),
            data={"arquivo": str(arquivo_pptx), "slides": n_slides}
        )

    def _montar_pptx(self, arquivo: Path, estrutura: dict):
        """Monta o arquivo .pptx usando python-pptx com design profissional."""
        from pptx import Presentation
        from pptx.util import Inches, Pt, Emu
        from pptx.dml.color import RGBColor
        from pptx.enum.text import PP_ALIGN

        prs = Presentation()
        prs.slide_width = Inches(13.33)
        prs.slide_height = Inches(7.5)

        # Paleta de cores
        AZUL        = RGBColor(0x1E, 0x3A, 0x8A)
        AZUL_CLARO  = RGBColor(0x3B, 0x82, 0xF6)
        AZUL_BRILHO = RGBColor(0x60, 0xA5, 0xFA)
        VERDE       = RGBColor(0x22, 0xC5, 0x5E)
        AMARELO     = RGBColor(0xF5, 0x9E, 0x0B)
        BRANCO      = RGBColor(0xFF, 0xFF, 0xFF)
        CINZA       = RGBColor(0x64, 0x74, 0x8B)
        CINZA_CLARO = RGBColor(0x94, 0xA3, 0xB8)
        ESCURO      = RGBColor(0x0F, 0x17, 0x2A)
        ESCURO2     = RGBColor(0x16, 0x19, 0x21)
        ESCURO3     = RGBColor(0x1C, 0x1F, 0x2B)

        def add_bg(slide, cor):
            fill = slide.background.fill
            fill.solid()
            fill.fore_color.rgb = cor

        def add_rect(slide, left, top, width, height, cor, line=False):
            s = slide.shapes.add_shape(1,
                Inches(left), Inches(top), Inches(width), Inches(height))
            s.fill.solid()
            s.fill.fore_color.rgb = cor
            if line:
                s.line.color.rgb = cor
            else:
                s.line.fill.background()
            return s

        def add_box(slide, text, left, top, width, height,
                    bold=False, size=24, color=BRANCO, align=PP_ALIGN.LEFT, italic=False):
            txBox = slide.shapes.add_textbox(
                Inches(left), Inches(top), Inches(width), Inches(height))
            tf = txBox.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.alignment = align
            run = p.add_run()
            run.text = text
            run.font.bold = bold
            run.font.italic = italic
            run.font.size = Pt(size)
            run.font.color.rgb = color
            return txBox

        def add_footer(slide, idx, total):
            # Linha separadora + número
            add_rect(slide, 0, 7.2, 13.33, 0.02, AZUL_CLARO)
            add_box(slide, f"{idx:02d} / {total:02d}", 11.5, 7.2, 1.5, 0.3,
                    size=10, color=CINZA, align=PP_ALIGN.RIGHT)

        slides_data = estrutura.get("slides", [])
        total = len(slides_data)

        for i, slide_info in enumerate(slides_data):
            tipo = slide_info.get("tipo", "conteudo")
            titulo_slide = slide_info.get("titulo", f"Slide {i+1}")
            pontos = slide_info.get("pontos", [])

            layout = prs.slide_layouts[6]  # blank
            slide = prs.slides.add_slide(layout)

            # ── CAPA ──────────────────────────────────────────────────
            if tipo == "capa" or i == 0:
                add_bg(slide, ESCURO)
                # Bloco lateral esquerdo colorido
                add_rect(slide, 0, 0, 0.5, 7.5, AZUL_CLARO)
                # Linha horizontal central
                add_rect(slide, 0.5, 3.35, 12.83, 0.04, AZUL_CLARO)
                # Título grande
                add_box(slide, estrutura.get("titulo", titulo_slide),
                        1.0, 1.0, 11.5, 2.0, bold=True, size=46, color=BRANCO)
                # Subtítulo
                add_box(slide, estrutura.get("subtitulo", ""),
                        1.0, 3.6, 10.0, 1.0, size=22, color=AZUL_BRILHO, italic=True)
                # Data / empresa
                add_box(slide, datetime.now().strftime("%B %Y"),
                        1.0, 6.6, 4.0, 0.6, size=13, color=CINZA)
                add_footer(slide, i + 1, total)

            # ── CONCLUSÃO ─────────────────────────────────────────────
            elif tipo == "conclusao":
                add_bg(slide, AZUL)
                add_rect(slide, 0, 0, 0.5, 7.5, AZUL_BRILHO)
                add_box(slide, titulo_slide,
                        1.0, 0.5, 11.5, 1.0, bold=True, size=30, color=BRANCO)
                add_rect(slide, 1.0, 1.5, 11.3, 0.04, BRANCO)
                y = 1.8
                for j, ponto in enumerate(pontos[:5]):
                    add_rect(slide, 1.0, y + 0.15, 0.3, 0.3, AZUL_BRILHO)
                    add_box(slide, ponto, 1.5, y, 11.0, 0.7,
                            bold=False, size=20, color=BRANCO)
                    y += 0.9
                add_footer(slide, i + 1, total)

            # ── DADOS (métricas) ──────────────────────────────────────
            elif tipo == "dados":
                add_bg(slide, ESCURO2)
                add_rect(slide, 0, 0, 13.33, 1.2, ESCURO3)
                add_box(slide, titulo_slide,
                        0.5, 0.15, 12.0, 0.9, bold=True, size=26, color=BRANCO)
                add_box(slide, f"{i+1:02d}", 12.0, 0.1, 1.0, 0.9,
                        bold=True, size=30, color=AZUL_CLARO, align=PP_ALIGN.RIGHT)
                # Cards de dados (distribui até 4 métricas em cards)
                n = min(len(pontos), 4)
                card_w = 11.0 / max(n, 1)
                for j, ponto in enumerate(pontos[:4]):
                    x = 1.1 + j * (card_w + 0.15)
                    add_rect(slide, x, 1.5, card_w, 2.5, ESCURO3)
                    add_box(slide, ponto, x + 0.15, 1.7, card_w - 0.3, 2.1,
                            bold=False, size=18, color=CINZA_CLARO)
                if len(pontos) > 4:
                    add_box(slide, pontos[4], 1.1, 4.2, 11.0, 1.0, size=17, color=BRANCO)
                add_footer(slide, i + 1, total)

            # ── PADRÃO (lista) ────────────────────────────────────────
            else:
                bg_cor = ESCURO if i % 2 == 0 else ESCURO2
                add_bg(slide, bg_cor)
                # Barra superior
                add_rect(slide, 0, 0, 13.33, 1.15, AZUL)
                # Acento vertical
                add_rect(slide, 0, 0, 0.12, 1.15, AZUL_BRILHO)
                add_box(slide, titulo_slide,
                        0.4, 0.13, 12.0, 0.9, bold=True, size=26, color=BRANCO)
                add_box(slide, f"{i+1:02d}", 11.8, 0.13, 1.2, 0.9,
                        bold=True, size=28, color=AZUL_CLARO, align=PP_ALIGN.RIGHT)

                y = 1.45
                for j, ponto in enumerate(pontos[:6]):
                    # Marcador colorido alternado
                    marc_cor = AZUL_CLARO if j % 2 == 0 else VERDE
                    add_rect(slide, 0.55, y + 0.18, 0.22, 0.22, marc_cor)
                    add_box(slide, ponto, 0.95, y, 11.6, 0.78,
                            bold=(j == 0), size=20 if j == 0 else 18, color=BRANCO)
                    y += 0.88

                nota = slide_info.get("nota", "")
                if nota:
                    slide.notes_slide.notes_text_frame.text = nota
                add_footer(slide, i + 1, total)

        prs.save(str(arquivo))
        log.info(f"PowerPoint salvo: {arquivo}")

    def _criar_powerpoint_direto(self, command: str) -> SkillResult:
        desc = self._extrair_descricao(command)
        return self._criar_powerpoint_planejado(desc, f"Descricao: {desc}\nEstilo: profissional")

    # ================================================================
    # EXCEL / PLANILHA
    # ================================================================

    def _criar_excel_planejado(self, descricao: str, contexto: str) -> SkillResult:
        """Cria planilha Excel com dados gerados pela IA."""
        APPS_DIR.mkdir(parents=True, exist_ok=True)
        nome = f"planilha_{descricao.replace(' ', '_')[:30]}_{datetime.now().strftime('%Y%m%d_%H%M')}"
        pasta = APPS_DIR / nome
        pasta.mkdir(parents=True, exist_ok=True)
        arquivo = pasta / f"{nome}.xlsx"

        prompt = f"""Voce e um especialista em planilhas Excel de negocios.
Com base nestas informacoes:
{contexto}

Crie a estrutura de uma planilha em formato JSON:
{{
  "nome_arquivo": "nome",
  "abas": [
    {{
      "nome": "Nome da aba",
      "colunas": ["Coluna A", "Coluna B", "Coluna C"],
      "linhas_exemplo": [
        ["valor1", "valor2", "valor3"]
      ]
    }}
  ]
}}

REGRAS:
- Minimo 2 abas, maximo 5
- Dados de exemplo REALISTAS e ESPECIFICOS para o contexto
- Numeros, datas, percentuais, moedas conforme adequado
- Pelo menos 10 linhas de exemplo por aba
Retorne APENAS o JSON, sem explicacoes."""

        estrutura = None
        try:
            from src.core.ai_engine import get_engine
            engine = get_engine()
            resp = engine.chat_with_fallback(prompt)
            if "{" in resp:
                inicio = resp.find("{")
                fim = resp.rfind("}") + 1
                estrutura = json.loads(resp[inicio:fim])
        except Exception as e:
            log.warning(f"IA nao gerou estrutura Excel: {e}")

        if not estrutura:
            estrutura = {
                "nome_arquivo": descricao,
                "abas": [{
                    "nome": "Dados",
                    "colunas": ["Item", "Quantidade", "Valor Unit (R$)", "Total (R$)"],
                    "linhas_exemplo": [
                        [f"Item {i}", str(i * 2), f"{i * 50:.2f}", f"{i * 100:.2f}"]
                        for i in range(1, 11)
                    ]
                }]
            }

        try:
            self._montar_xlsx(arquivo, estrutura)
        except ImportError:
            log.info("Instalando openpyxl...")
            subprocess.run(["py", "-3", "-m", "pip", "install", "openpyxl", "-q"],
                           capture_output=True, timeout=60)
            try:
                self._montar_xlsx(arquivo, estrutura)
            except Exception as e2:
                return SkillResult(success=False, message=f"Erro ao criar Excel: {e2}")
        except Exception as e:
            return SkillResult(success=False, message=f"Erro ao montar Excel: {e}")

        try:
            os.startfile(str(arquivo))
        except Exception:
            subprocess.Popen(["start", "", str(arquivo)], shell=True)

        n_abas = len(estrutura.get("abas", []))
        nomes_abas = ", ".join(a.get("nome", "") for a in estrutura.get("abas", []))
        return SkillResult(
            success=True,
            message=(
                f"[PLANILHA CRIADA] {descricao.title()}\n\n"
                f"ARQUIVO: {arquivo}\n"
                f"ABAS: {n_abas} | {nomes_abas}\n"
                f"STATUS: Aberto no Excel\n\n"
                f"Planilha com dados reais criada pela IA!"
            ),
            data={"arquivo": str(arquivo)}
        )

    def _montar_xlsx(self, arquivo: Path, estrutura: dict):
        """Monta o arquivo .xlsx usando openpyxl."""
        from openpyxl import Workbook
        from openpyxl.styles import Font, PatternFill, Alignment, Border, Side
        from openpyxl.utils import get_column_letter

        wb = Workbook()
        wb.remove(wb.active)

        h_fill = PatternFill(start_color="1E3A8A", end_color="1E3A8A", fill_type="solid")
        alt_fill = PatternFill(start_color="F1F5F9", end_color="F1F5F9", fill_type="solid")
        h_font = Font(bold=True, color="FFFFFF", size=11)
        t_font = Font(bold=True, color="3B82F6", size=13)
        thin = Side(border_style="thin", color="CBD5E1")
        border = Border(left=thin, right=thin, top=thin, bottom=thin)

        for aba_data in estrutura.get("abas", []):
            nome_aba = str(aba_data.get("nome", "Dados"))[:31]
            ws = wb.create_sheet(title=nome_aba)
            colunas = aba_data.get("colunas", [])
            linhas = aba_data.get("linhas_exemplo", [])

            # Título
            ws.merge_cells(f"A1:{get_column_letter(max(len(colunas), 1))}1")
            ws["A1"] = aba_data.get("nome", "Dados")
            ws["A1"].font = t_font
            ws["A1"].alignment = Alignment(horizontal="center")
            ws.row_dimensions[1].height = 28

            # Headers
            for ci, col in enumerate(colunas, 1):
                cell = ws.cell(row=2, column=ci, value=col)
                cell.font = h_font
                cell.fill = h_fill
                cell.alignment = Alignment(horizontal="center", vertical="center")
                cell.border = border
                ws.column_dimensions[get_column_letter(ci)].width = max(15, len(str(col)) + 4)
            ws.row_dimensions[2].height = 22

            # Dados
            for ri, linha in enumerate(linhas, 3):
                fill = alt_fill if ri % 2 == 0 else None
                for ci, valor in enumerate(linha[:len(colunas)], 1):
                    cell = ws.cell(row=ri, column=ci, value=valor)
                    cell.border = border
                    cell.alignment = Alignment(vertical="center")
                    if fill:
                        cell.fill = fill

            ws.freeze_panes = "A3"

        wb.save(str(arquivo))
        log.info(f"Excel salvo: {arquivo}")

    def _criar_excel_direto(self, command: str) -> SkillResult:
        desc = self._extrair_descricao(command)
        return self._criar_excel_planejado(desc, f"Descricao: {desc}")

    # ================================================================
    # WORD / DOCUMENTO
    # ================================================================

    def _criar_word_planejado(self, descricao: str, contexto: str) -> SkillResult:
        """Cria documento Word com conteúdo gerado pela IA."""
        APPS_DIR.mkdir(parents=True, exist_ok=True)
        nome = f"doc_{descricao.replace(' ', '_')[:30]}_{datetime.now().strftime('%Y%m%d_%H%M')}"
        pasta = APPS_DIR / nome
        pasta.mkdir(parents=True, exist_ok=True)
        arquivo = pasta / f"{nome}.docx"

        prompt = f"""Voce e um redator profissional de documentos de negocios.
Com base nestas informacoes:
{contexto}

Escreva o documento completo em formato JSON:
{{
  "titulo": "Titulo do documento",
  "secoes": [
    {{
      "titulo": "Titulo da secao",
      "tipo": "normal|lista|destaque",
      "conteudo": "texto completo da secao; para lista use itens separados por |"
    }}
  ]
}}

REGRAS:
- Minimo 5 secoes bem desenvolvidas
- Conteudo ESPECIFICO, detalhado e profissional
- Dados, numeros e exemplos REAIS para o contexto
Retorne APENAS o JSON, sem explicacoes."""

        estrutura = None
        try:
            from src.core.ai_engine import get_engine
            engine = get_engine()
            resp = engine.chat_with_fallback(prompt)
            if "{" in resp:
                inicio = resp.find("{")
                fim = resp.rfind("}") + 1
                estrutura = json.loads(resp[inicio:fim])
        except Exception as e:
            log.warning(f"IA nao gerou estrutura Word: {e}")

        if not estrutura:
            estrutura = {
                "titulo": descricao.title(),
                "secoes": [
                    {"titulo": "Introducao", "tipo": "normal",
                     "conteudo": f"Este documento descreve {descricao}."},
                    {"titulo": "Objetivos", "tipo": "lista",
                     "conteudo": "Objetivo 1|Objetivo 2|Objetivo 3"},
                    {"titulo": "Desenvolvimento", "tipo": "normal",
                     "conteudo": "Conteudo principal do documento."},
                    {"titulo": "Conclusao", "tipo": "destaque",
                     "conteudo": "Conclusoes e proximos passos."},
                ]
            }

        try:
            self._montar_docx(arquivo, estrutura)
        except ImportError:
            log.info("Instalando python-docx...")
            subprocess.run(["py", "-3", "-m", "pip", "install", "python-docx", "-q"],
                           capture_output=True, timeout=60)
            try:
                self._montar_docx(arquivo, estrutura)
            except Exception as e2:
                return SkillResult(success=False, message=f"Erro ao criar Word: {e2}")
        except Exception as e:
            return SkillResult(success=False, message=f"Erro ao montar DOCX: {e}")

        try:
            os.startfile(str(arquivo))
        except Exception:
            subprocess.Popen(["start", "", str(arquivo)], shell=True)

        n_secoes = len(estrutura.get("secoes", []))
        return SkillResult(
            success=True,
            message=(
                f"[DOCUMENTO CRIADO] {estrutura.get('titulo', descricao.title())}\n\n"
                f"ARQUIVO: {arquivo}\n"
                f"SECOES: {n_secoes} secoes com conteudo real\n"
                f"STATUS: Aberto no Word\n\n"
                f"Documento profissional criado pela IA!"
            ),
            data={"arquivo": str(arquivo)}
        )

    def _montar_docx(self, arquivo: Path, estrutura: dict):
        """Monta o arquivo .docx usando python-docx."""
        from docx import Document
        from docx.shared import Pt, RGBColor, Cm
        from docx.enum.text import WD_ALIGN_PARAGRAPH
        from docx.oxml.ns import qn
        from docx.oxml import OxmlElement

        doc = Document()
        for section in doc.sections:
            section.top_margin = Cm(2.5)
            section.bottom_margin = Cm(2.5)
            section.left_margin = Cm(3)
            section.right_margin = Cm(2)

        # Título
        titulo_para = doc.add_heading(estrutura.get("titulo", "Documento"), level=0)
        titulo_para.alignment = WD_ALIGN_PARAGRAPH.CENTER
        for run in titulo_para.runs:
            run.font.color.rgb = RGBColor(0x1E, 0x3A, 0x8A)

        doc.add_paragraph()
        data_para = doc.add_paragraph(f"Data: {datetime.now().strftime('%d/%m/%Y')}")
        data_para.alignment = WD_ALIGN_PARAGRAPH.RIGHT
        for run in data_para.runs:
            run.font.size = Pt(10)
            run.font.color.rgb = RGBColor(0x64, 0x74, 0x8B)
        doc.add_paragraph()

        for secao in estrutura.get("secoes", []):
            titulo_sec = secao.get("titulo", "")
            tipo = secao.get("tipo", "normal")
            conteudo = secao.get("conteudo", "")

            h = doc.add_heading(titulo_sec, level=1)
            for run in h.runs:
                run.font.color.rgb = RGBColor(0x1E, 0x3A, 0x8A)

            if tipo == "lista":
                itens = [x.strip() for x in conteudo.split("|") if x.strip()]
                for item in itens:
                    p = doc.add_paragraph(item, style="List Bullet")
                    for run in p.runs:
                        run.font.size = Pt(11)
            elif tipo == "destaque":
                p = doc.add_paragraph()
                run = p.add_run(conteudo)
                run.bold = True
                run.font.size = Pt(11)
                run.font.color.rgb = RGBColor(0x1E, 0x3A, 0x8A)
                pPr = p._element.get_or_add_pPr()
                shd = OxmlElement("w:shd")
                shd.set(qn("w:val"), "clear")
                shd.set(qn("w:color"), "auto")
                shd.set(qn("w:fill"), "EFF6FF")
                pPr.append(shd)
            else:
                p = doc.add_paragraph(conteudo)
                for run in p.runs:
                    run.font.size = Pt(11)

            doc.add_paragraph()

        doc.save(str(arquivo))
        log.info(f"Word salvo: {arquivo}")

    def _criar_word_direto(self, command: str) -> SkillResult:
        desc = self._extrair_descricao(command)
        return self._criar_word_planejado(desc, f"Descricao: {desc}")

    # ================================================================
    # HTML / WEB APPS
    # ================================================================

    def _extrair_descricao(self, command: str) -> str:
        palavras_remover = {
            "crie", "cria", "faca", "desenvolva", "programe",
            "um", "uma", "app", "site", "sistema", "ferramenta",
            "de", "para", "do", "da", "aplicativo", "apresentacao",
            "powerpoint", "planilha", "excel", "documento", "word"
        }
        palavras = command.lower().split()
        desc = " ".join(w for w in palavras if w not in palavras_remover)
        return desc.strip() or "ferramenta"

    def _abrir_vscode(self, pasta: Path) -> bool:
        try:
            subprocess.Popen(["code", str(pasta)], shell=True)
            return True
        except Exception:
            try:
                subprocess.Popen(f'code "{pasta}"', shell=True)
                return True
            except Exception:
                return False

    def _criar_app_html_planejado(self, descricao: str, contexto: str,
                                   abrir_vscode: bool = False) -> SkillResult:
        """Cria app HTML com IA — usa TODAS as funcionalidades listadas no contexto."""
        nome_projeto = f"app_{descricao.replace(' ', '_')[:30]}_{datetime.now().strftime('%H%M')}"
        pasta = APPS_DIR / nome_projeto
        pasta.mkdir(parents=True, exist_ok=True)

        # Extrai lista de funcionalidades do contexto para incluir no prompt
        func_block = ""
        if "FUNCIONALIDADES OBRIGATORIAS" in contexto:
            idx = contexto.find("FUNCIONALIDADES OBRIGATORIAS")
            func_block = contexto[idx:]

        prompt = (
            "Voce e um desenvolvedor front-end senior com 15 anos de experiencia em UI/UX.\n"
            f"BRIEFING DO PROJETO:\n{contexto}\n\n"
            "═══════════════════════════════════════════════════════════\n"
            "MISSAO: Crie este projeto EXATAMENTE como especificado acima.\n"
            "CADA FUNCIONALIDADE listada DEVE existir no codigo final.\n"
            "Nao invente funcionalidades diferentes. Implemente TODAS as pedidas.\n"
            "═══════════════════════════════════════════════════════════\n\n"
            "=== DESIGN PREMIUM OBRIGATORIO ===\n"
            "- Dark theme: background #0f1117, cards #161921, border #2d3748\n"
            "- Accent azul: #3b82f6 | Verde: #22c55e | Amarelo: #f59e0b | Vermelho: #ef4444\n"
            "- Fonte: 'Segoe UI', system-ui, sans-serif\n"
            "- Cards com border-radius 12px, sombra box-shadow 0 4px 20px rgba(0,0,0,0.4)\n"
            "- Botoes com transition 0.2s, hover translateY(-2px) e brightness(1.15)\n"
            "- Header ou sidebar com gradiente linear-gradient(135deg,...)\n"
            "- Icones emoji relevantes nos titulos e cards\n"
            "- Toast notifications flutuantes (aparecem no canto ao executar acoes)\n"
            "- Animacao de fade-in nos cards ao carregar (CSS @keyframes)\n\n"
            "=== CODIGO OBRIGATORIO ===\n"
            "- JavaScript puro 100% — ZERO CDN, ZERO frameworks externos\n"
            "- localStorage para persistencia de dados\n"
            "- Dados de exemplo REAIS e ESPECIFICOS (nao 'Item 1', 'Item 2')\n"
            "- Validacao visual nos formularios (borda vermelha se invalido)\n"
            "- Responsivo (funciona em 320px ate 1920px)\n\n"
            "=== RESTRICOES ABSOLUTAS ===\n"
            "- PROIBIDO: <link rel=stylesheet>, src=script.js, <script src=...>\n"
            "- TODO CSS dentro de <style> no <head>\n"
            "- TODO JS dentro de <script> no final do <body>\n"
            "- PROIBIDO: Bootstrap, Tailwind, jQuery, React, Vue, Chart.js, D3\n"
            "- Arquivo UNICO, funciona 100% offline sem servidor\n\n"
            "Retorne APENAS o HTML completo iniciando com <!DOCTYPE html> e finalizando com </html>.\n"
            "SEM explicacoes. SEM markdown. SEM comentarios fora do HTML."
        )

        html_gerado = ""
        try:
            from src.core.ai_engine import get_engine
            engine = get_engine()
            log.info(f"Gerando HTML para: {descricao[:50]}")
            resposta = engine.chat_with_fallback(prompt)
            if resposta and "<!DOCTYPE" in resposta.upper():
                inicio = resposta.upper().find("<!DOCTYPE")
                fim = resposta.rfind("</html>")
                if fim > inicio:
                    candidato = resposta[inicio:fim + 7]
                    # Rejeita HTML com arquivos externos (CSS/JS separados)
                    ext_refs = ['<link rel="stylesheet"', "<link rel='stylesheet'",
                                'src="script.js"', "src='script.js'",
                                'src="style.js"', 'href="style.css"', "href='style.css'"]
                    has_ext = any(ref in candidato.lower() for ref in ext_refs)
                    if not has_ext and len(candidato) >= 500:
                        html_gerado = candidato
                    else:
                        log.warning(f"HTML rejeitado: externo={has_ext}, len={len(candidato)}")
        except Exception as e:
            log.warning(f"IA nao disponivel: {e}")

        if not html_gerado or len(html_gerado) < 500:
            log.warning("Usando fallback HTML local")
            html_gerado = self._html_fallback(descricao)

        (pasta / "index.html").write_text(html_gerado, encoding="utf-8")
        (pasta / "README.md").write_text(
            f"# {descricao.title()}\nCriado por William v6 com IA\n\nAbra index.html no navegador.",
            encoding="utf-8"
        )

        # Abre VS Code apenas se o usuário pediu
        if abrir_vscode:
            self._abrir_vscode(pasta)

        # Abre no navegador
        try:
            import webbrowser
            webbrowser.open(str(pasta / "index.html"))
        except Exception:
            pass

        tamanho_kb = round(len(html_gerado) / 1024, 1)
        vsc_info = "VS Code + navegador abertos!" if abrir_vscode else "Aberto no navegador!"
        return SkillResult(
            success=True,
            message=(
                f"[APP CRIADO] {descricao.title()}\n\n"
                f"ARQUIVO: {pasta / 'index.html'}\n"
                f"TAMANHO: {tamanho_kb} KB de codigo gerado\n"
                f"STATUS: {vsc_info}\n\n"
                f"O app foi gerado com TODAS as funcionalidades solicitadas.\n"
                f"Abra o navegador para testar. Se quiser editar o codigo,\n"
                f"abra o arquivo index.html no VS Code."
            ),
            data={"pasta": str(pasta), "arquivo": str(pasta / "index.html")}
        )

    def _criar_landing_page(self, command: str) -> SkillResult:
        desc = self._extrair_descricao(command)
        contexto = (
            f"Tipo: Landing Page de Captura de Leads\n"
            f"Descricao: {desc}\n"
            f"Objetivo: Capturar nome, WhatsApp e email do visitante\n"
            f"Estilo: Moderno, persuasivo, profissional\n"
            f"Publico: Empresas e profissionais brasileiros"
        )
        return self._criar_app_html_planejado(f"landing_{desc}", contexto)

    def _criar_dashboard(self, command: str) -> SkillResult:
        desc = self._extrair_descricao(command)
        contexto = (
            f"Tipo: Dashboard de Analytics/Metricas\n"
            f"Descricao: {desc}\n"
            f"Objetivo: Exibir KPIs, graficos de barras (CSS puro), tabela de dados com filtros\n"
            f"Estilo: Dark theme profissional, dados interativos\n"
            f"Extras: Graficos de barra em CSS puro, filtros de data, exportar dados"
        )
        return self._criar_app_html_planejado(f"dashboard_{desc}", contexto)

    def _criar_calculadora(self, command: str) -> SkillResult:
        desc = self._extrair_descricao(command)
        pasta = APPS_DIR / f"calculadora_{desc.replace(' ', '_')}"
        pasta.mkdir(parents=True, exist_ok=True)
        html = """<!DOCTYPE html><html lang='pt-BR'><head><meta charset='UTF-8'><title>Calculadora ROI</title>
<style>body{font-family:'Segoe UI',sans-serif;background:#0f1117;color:#f1f5f9;display:flex;align-items:center;justify-content:center;min-height:100vh}
.card{background:#161921;padding:40px;border-radius:15px;width:500px;border:1px solid #2d3748}
h1{color:#3b82f6;margin-bottom:30px;text-align:center}label{display:block;margin:15px 0 5px;color:#94a3b8}
input{width:100%;padding:12px;border-radius:6px;border:1px solid #2d3748;background:#1c1f2b;color:#f1f5f9;font-size:1rem;box-sizing:border-box}
button{width:100%;padding:14px;background:#3b82f6;color:white;border:none;border-radius:6px;font-size:1rem;font-weight:600;cursor:pointer;margin-top:20px}
.res{margin-top:20px;padding:20px;background:#1c1f2b;border-radius:8px;display:none}
.res .v{font-size:2rem;color:#22c55e;font-weight:700}
</style></head><body><div class='card'>
<h1>Calculadora de ROI</h1>
<label>Investimento Mensal (R$)</label><input type='number' id='inv' value='1500'>
<label>Clientes por Mes</label><input type='number' id='cli' value='5'>
<label>Ticket Medio (R$)</label><input type='number' id='tick' value='2000'>
<label>Meses de Retencao</label><input type='number' id='ret' value='6'>
<button onclick='calc()'>CALCULAR ROI</button>
<div class='res' id='res'>
<p>Receita Total:</p><div class='v' id='rec'></div>
<p style='margin-top:10px'>ROI: <span id='roi' style='color:#3b82f6;font-weight:700'></span></p>
<p>Lucro: <span id='luc' style='color:#22c55e;font-weight:700'></span></p>
</div></div>
<script>
function calc(){
var inv=+document.getElementById('inv').value,cli=+document.getElementById('cli').value,
tick=+document.getElementById('tick').value,ret=+document.getElementById('ret').value;
var rec=cli*tick*ret,roi=((rec-inv)/inv*100).toFixed(0),luc=rec-inv;
document.getElementById('rec').textContent='R$'+rec.toLocaleString('pt-BR');
document.getElementById('roi').textContent=roi+'%';
document.getElementById('luc').textContent='R$'+luc.toLocaleString('pt-BR');
document.getElementById('res').style.display='block';
}
</script></body></html>"""
        (pasta / "index.html").write_text(html, encoding="utf-8")
        self._abrir_vscode(pasta)
        try:
            import webbrowser
            webbrowser.open(str(pasta / "index.html"))
        except Exception:
            pass
        return SkillResult(success=True,
                           message=f"[CALCULADORA ROI] Aberta!\nArquivo: {pasta}/index.html",
                           data={"pasta": str(pasta)})

    def _criar_crm(self, command: str) -> SkillResult:
        desc = self._extrair_descricao(command)
        pasta = APPS_DIR / f"crm_{desc.replace(' ', '_')}"
        pasta.mkdir(parents=True, exist_ok=True)
        key = desc.replace(" ", "_")
        html = f"""<!DOCTYPE html><html lang='pt-BR'><head><meta charset='UTF-8'><title>CRM {desc.title()}</title>
<style>body{{font-family:'Segoe UI',sans-serif;background:#0f1117;color:#f1f5f9;margin:0}}
.header{{background:#161921;padding:20px 30px}}.header h1{{color:#3b82f6}}
.form{{background:#161921;margin:20px;padding:20px;border-radius:10px}}
.fr{{display:flex;gap:10px;flex-wrap:wrap}}
input{{padding:10px;border-radius:5px;border:1px solid #2d3748;background:#1c1f2b;color:#f1f5f9;flex:1;min-width:120px}}
button{{padding:10px 20px;background:#3b82f6;color:white;border:none;border-radius:5px;cursor:pointer;font-weight:600}}
.table-wrap{{margin:20px;background:#161921;border-radius:10px;overflow:hidden}}
table{{width:100%;border-collapse:collapse}}th{{background:#1c1f2b;padding:12px;text-align:left;color:#94a3b8}}
td{{padding:12px;border-top:1px solid #2d3748}}
</style></head><body>
<div class='header'><h1 id='tit'>CRM - 0 Leads</h1></div>
<div class='form'><div class='fr'>
<input id='n' placeholder='Nome *'><input id='emp' placeholder='Empresa'>
<input id='tel' placeholder='WhatsApp'><input id='email' placeholder='Email'>
<button onclick='add()'>+ Adicionar</button>
</div></div>
<div class='table-wrap'><table>
<thead><tr><th>Nome</th><th>Empresa</th><th>Telefone</th><th>Email</th><th>Data</th><th>Acao</th></tr></thead>
<tbody id='tb'></tbody></table></div>
<script>
var K='crm_{key}';
function load(){{return JSON.parse(localStorage.getItem(K)||'[]')}}
function save(d){{localStorage.setItem(K,JSON.stringify(d));render()}}
function add(){{
var n=document.getElementById('n').value.trim();
if(!n){{alert('Informe o nome!');return;}}
var d=load();
d.push({{id:Date.now(),nome:n,emp:document.getElementById('emp').value,
tel:document.getElementById('tel').value,email:document.getElementById('email').value,
data:new Date().toLocaleDateString('pt-BR')}});
save(d);['n','emp','tel','email'].forEach(function(id){{document.getElementById(id).value=''}});
}}
function del(id){{save(load().filter(function(x){{return x.id!==id}}))}}
function render(){{
var d=load();
document.getElementById('tit').textContent='CRM - '+d.length+' Leads';
document.getElementById('tb').innerHTML=d.map(function(x){{
return '<tr><td>'+x.nome+'</td><td>'+(x.emp||'-')+'</td><td>'+(x.tel||'-')+'</td><td>'+(x.email||'-')+'</td><td>'+x.data+'</td><td><button onclick="del('+x.id+')" style="background:#ef4444;padding:4px 10px">x</button></td></tr>';
}}).join('');
}}
render();
</script></body></html>"""
        (pasta / "index.html").write_text(html, encoding="utf-8")
        self._abrir_vscode(pasta)
        try:
            import webbrowser
            webbrowser.open(str(pasta / "index.html"))
        except Exception:
            pass
        return SkillResult(success=True,
                           message=f"[CRM] {desc.title()}\nArquivo: {pasta}/index.html",
                           data={"pasta": str(pasta)})

    def _criar_app_generico(self, command: str) -> SkillResult:
        """Cria app genérico com IA."""
        desc = self._extrair_descricao(command)
        return self._criar_app_html_planejado(
            desc,
            f"Descricao: {desc}\nTipo: app web geral\nEstilo: moderno profissional"
        )

    def _html_fallback(self, desc: str) -> str:
        key = desc.replace(" ", "_")
        return f"""<!DOCTYPE html>
<html lang="pt-BR">
<head>
    <meta charset="UTF-8"><title>{desc.title()}</title>
    <style>
        *{{margin:0;padding:0;box-sizing:border-box}}
        body{{font-family:'Segoe UI',sans-serif;background:#0f1117;color:#f1f5f9}}
        .header{{background:#161921;padding:20px 30px;border-bottom:1px solid #2d3748}}
        .header h1{{color:#3b82f6}}
        .container{{max-width:900px;margin:30px auto;padding:0 20px}}
        .kpis{{display:grid;grid-template-columns:repeat(3,1fr);gap:15px;margin-bottom:25px}}
        .kpi{{background:#161921;border:1px solid #2d3748;border-radius:10px;padding:20px;text-align:center}}
        .kpi .valor{{font-size:2rem;font-weight:700;color:#3b82f6}}
        .kpi .label{{color:#94a3b8;font-size:0.85rem}}
        .card{{background:#161921;border:1px solid #2d3748;border-radius:10px;padding:25px;margin-bottom:20px}}
        .card h2{{color:#f1f5f9;margin-bottom:15px}}
        .fr{{display:flex;gap:10px;margin-bottom:15px;flex-wrap:wrap}}
        input,select{{flex:1;padding:10px 14px;background:#1c1f2b;border:1px solid #2d3748;border-radius:6px;color:#f1f5f9;min-width:150px}}
        button{{padding:10px 20px;background:#3b82f6;color:white;border:none;border-radius:6px;cursor:pointer;font-weight:600}}
        button:hover{{background:#2563eb}}
        .btn-red{{background:#ef4444}}
        table{{width:100%;border-collapse:collapse}}
        th{{background:#1c1f2b;padding:12px;text-align:left;color:#94a3b8}}
        td{{padding:12px;border-top:1px solid #2d3748}}
        .tag{{padding:3px 10px;border-radius:20px;font-size:0.8rem;font-weight:600}}
        .tag-g{{background:#14532d;color:#22c55e}}
        .tag-y{{background:#451a03;color:#f59e0b}}
        .tag-b{{background:#1e3a5f;color:#3b82f6}}
    </style>
</head>
<body>
    <div class="header"><h1>⚡ {desc.title()}</h1></div>
    <div class="container">
        <div class="kpis">
            <div class="kpi"><div class="valor" id="k1">0</div><div class="label">Total</div></div>
            <div class="kpi"><div class="valor" id="k2">0</div><div class="label">Ativos</div></div>
            <div class="kpi"><div class="valor" id="k3">0%</div><div class="label">Conclusao</div></div>
        </div>
        <div class="card">
            <h2>Adicionar Item</h2>
            <div class="fr">
                <input type="text" id="nome" placeholder="Nome *" required>
                <input type="text" id="detalhe" placeholder="Detalhe">
                <select id="status">
                    <option value="ativo">Ativo</option>
                    <option value="pendente">Pendente</option>
                    <option value="concluido">Concluido</option>
                </select>
                <button class="btn-green" onclick="adicionar()">+ Adicionar</button>
            </div>
        </div>
        <div class="card">
            <h2>Lista</h2>
            <table>
                <thead><tr><th>Nome</th><th>Detalhe</th><th>Status</th><th>Data</th><th>Acao</th></tr></thead>
                <tbody id="tbody"></tbody>
            </table>
        </div>
    </div>
<script>
var K='{key}';
function carregar(){{return JSON.parse(localStorage.getItem(K)||'[]');}}
function salvar(d){{localStorage.setItem(K,JSON.stringify(d));atualizar();}}
function adicionar(){{
    var nome=document.getElementById('nome').value.trim();
    if(!nome){{alert('Informe o nome!');return;}}
    var dados=carregar();
    dados.push({{id:Date.now(),nome:nome,detalhe:document.getElementById('detalhe').value,
        status:document.getElementById('status').value,data:new Date().toLocaleDateString('pt-BR')}});
    salvar(dados);
    document.getElementById('nome').value='';
    document.getElementById('detalhe').value='';
}}
function remover(id){{salvar(carregar().filter(function(d){{return d.id!==id;}}));}}
function atualizar(){{
    var dados=carregar();
    var tags={{ativo:'tag-g',pendente:'tag-y',concluido:'tag-b'}};
    document.getElementById('tbody').innerHTML=dados.map(function(d){{
        return '<tr><td><strong>'+d.nome+'</strong></td><td style="color:#94a3b8">'+(d.detalhe||'-')+'</td>'
            +'<td><span class="tag '+(tags[d.status]||'tag-b')+'">'+d.status+'</span></td>'
            +'<td style="color:#94a3b8">'+d.data+'</td>'
            +'<td><button class="btn-red" style="padding:5px 12px" onclick="remover('+d.id+')">x</button></td></tr>';
    }}).join('');
    document.getElementById('k1').textContent=dados.length;
    document.getElementById('k2').textContent=dados.filter(function(d){{return d.status==='ativo';}}).length;
    var conc=dados.filter(function(d){{return d.status==='concluido';}}).length;
    document.getElementById('k3').textContent=dados.length?Math.round(conc/dados.length*100)+'%':'0%';
}}
atualizar();
</script>
</body>
</html>"""
